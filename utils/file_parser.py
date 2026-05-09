import os
import logging
from typing import Optional

try:
    import pypdf
except ImportError:
    pypdf = None
try:
    import docx
except ImportError:
    docx = None
try:
    import openpyxl
except ImportError:
    openpyxl = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

logger = logging.getLogger(__name__)

def extract_text_from_file(file_path: str, file_name: str) -> Optional[str]:
    ext = os.path.splitext(file_name)[1].lower()
    
    try:
        if ext == '.pdf' and pypdf:
            text = ""
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text
        elif ext == '.docx' and docx:
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        elif ext == '.xlsx' and openpyxl:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text = ""
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
            return text
        elif ext == '.pptx' and Presentation:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        elif ext in ['.html', '.htm'] and BeautifulSoup:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f, 'html.parser')
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text(separator='\n')
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return None
    except Exception as e:
        logger.error(f"Error extracting text from {file_name}: {e}")
        return None